#!/usr/bin/env python3
"""Build corporate-styled DOCX and PPTX from Markdown sources.

Generates properly styled Word documents with:
- Cover page with title/author/date
- Table of contents
- Styled tables with borders and header highlighting
- Code blocks with monospace font
- Mermaid diagrams as descriptive placeholders
- Page numbers in footer

Usage:
    python scripts/build_corporate.py --word     # Generate styled DOCX
    python scripts/build_corporate.py --slides   # Generate PPTX
    python scripts/build_corporate.py --all      # Both
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from docx import Document
from docx.shared import Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Cm as PptxCm
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "output"
LIVRABLES_DIR = ROOT / "docs" / "04-livrables"

LIVRABLE_FILES = [
    LIVRABLES_DIR / "architecture-technique.md",
    LIVRABLES_DIR / "strategie-migration.md",
    LIVRABLES_DIR / "guide-depannage-toip.md",
    LIVRABLES_DIR / "note-recommandation-wms.md",
    LIVRABLES_DIR / "configs" / "vpn-ipsec.md",
    LIVRABLES_DIR / "configs" / "pare-feu.md",
]

# Corporate colors
BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_MED = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
BLUE_PALE = RGBColor(0xE8, 0xF0, 0xF8)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MED = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)


# ============================================================
# MARKDOWN PARSER
# ============================================================

def parse_frontmatter(text: str) -> tuple[dict, str]:
    """Extract YAML frontmatter and return (metadata, body)."""
    if not text.startswith("---"):
        return {}, text
    end = text.find("---", 3)
    if end == -1:
        return {}, text
    yaml_block = text[3:end].strip()
    body = text[end + 3:].strip()
    meta = {}
    for line in yaml_block.split("\n"):
        if ":" in line:
            key, val = line.split(":", 1)
            meta[key.strip()] = val.strip().strip('"').strip("'")
    return meta, body


def parse_table(lines: list[str]) -> list[list[str]]:
    """Parse markdown pipe table lines into rows of cells."""
    rows = []
    for line in lines:
        line = line.strip()
        if line.startswith("|"):
            line = line[1:]
        if line.endswith("|"):
            line = line[:-1]
        cells = [c.strip() for c in line.split("|")]
        # Skip separator rows (|---|---|)
        if all(re.match(r'^[-:]+$', c) for c in cells):
            continue
        rows.append(cells)
    return rows


def parse_markdown_blocks(text: str) -> list[dict]:
    """Parse markdown into a list of blocks."""
    blocks = []
    lines = text.split("\n")
    i = 0

    while i < len(lines):
        line = lines[i]

        # Blank line
        if not line.strip():
            i += 1
            continue

        # Horizontal rule
        if re.match(r'^---+\s*$', line.strip()):
            blocks.append({"type": "hr"})
            i += 1
            continue

        # Heading
        m = re.match(r'^(#{1,6})\s+(.+)', line)
        if m:
            blocks.append({
                "type": "heading",
                "level": len(m.group(1)),
                "text": m.group(2).strip()
            })
            i += 1
            continue

        # Code block (fenced)
        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            if lang == "mermaid":
                # Extract a description from the mermaid code
                desc = _describe_mermaid(code_lines)
                blocks.append({"type": "mermaid", "description": desc, "code": "\n".join(code_lines)})
            else:
                blocks.append({"type": "code", "lang": lang, "text": "\n".join(code_lines)})
            continue

        # Table
        if "|" in line and i + 1 < len(lines) and "|" in lines[i + 1]:
            table_lines = []
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                table_lines.append(lines[i])
                i += 1
            rows = parse_table(table_lines)
            if rows:
                blocks.append({"type": "table", "rows": rows})
            continue

        # Blockquote
        if line.strip().startswith(">"):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(lines[i].strip().lstrip(">").strip())
                i += 1
            blocks.append({"type": "blockquote", "text": " ".join(quote_lines)})
            continue

        # Unordered list
        if re.match(r'^[-*]\s+', line.strip()):
            items = []
            while i < len(lines) and re.match(r'^[-*]\s+', lines[i].strip()):
                items.append(lines[i].strip().lstrip("-*").strip())
                i += 1
            blocks.append({"type": "list", "items": items})
            continue

        # Numbered list
        if re.match(r'^\d+\.\s+', line.strip()):
            items = []
            while i < len(lines) and re.match(r'^\d+\.\s+', lines[i].strip()):
                items.append(re.sub(r'^\d+\.\s+', '', lines[i].strip()))
                i += 1
            blocks.append({"type": "numbered_list", "items": items})
            continue

        # Regular paragraph
        para_lines = []
        while i < len(lines) and lines[i].strip() and not lines[i].startswith("#") \
                and not lines[i].strip().startswith("```") and not lines[i].strip().startswith("|") \
                and not lines[i].strip().startswith(">") and not re.match(r'^---+\s*$', lines[i].strip()) \
                and not re.match(r'^[-*]\s+', lines[i].strip()) \
                and not re.match(r'^\d+\.\s+', lines[i].strip()):
            para_lines.append(lines[i].strip())
            i += 1
        if para_lines:
            blocks.append({"type": "paragraph", "text": " ".join(para_lines)})

    return blocks


def _describe_mermaid(code_lines: list[str]) -> str:
    """Generate a text description of a mermaid diagram."""
    code = "\n".join(code_lines)
    subgraphs = re.findall(r'subgraph\s+\w+\["([^"]+)"\]', code)
    if subgraphs:
        return "Diagramme reseau : " + " | ".join(subgraphs)
    return "Schema d'architecture (voir source Markdown pour le detail)"


# ============================================================
# DOCX BUILDER
# ============================================================

def _add_formatted_runs(paragraph, text: str) -> None:
    """Add runs with bold/italic/code formatting from markdown inline syntax."""
    # Pattern for **bold**, *italic*, `code`
    parts = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)', text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run(part[2:-2])
            run.bold = True
        elif part.startswith("*") and part.endswith("*") and not part.startswith("**"):
            run = paragraph.add_run(part[1:-1])
            run.italic = True
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(9)
            run.font.color.rgb = BLUE_MED
        else:
            paragraph.add_run(part)


def _set_cell_shading(cell, color_hex: str) -> None:
    """Set background color on a table cell."""
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def _set_table_borders(table) -> None:
    """Add borders to all cells in a table."""
    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else parse_xml(f'<w:tblPr {nsdecls("w")}/>')
    borders = parse_xml(
        f'<w:tblBorders {nsdecls("w")}>'
        '  <w:top w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '  <w:left w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '  <w:bottom w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '  <w:right w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '  <w:insideH w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '  <w:insideV w:val="single" w:sz="4" w:space="0" w:color="BFBFBF"/>'
        '</w:tblBorders>'
    )
    tblPr.append(borders)


def _add_page_number(doc: Document) -> None:
    """Add page number to footer."""
    section = doc.sections[-1]
    footer = section.footer
    footer.is_linked_to_previous = False
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Company name on left
    run = p.add_run("MSPR - NordTransit Logistics")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY_MED

    run = p.add_run("    |    Page ")
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY_MED

    # Page number field
    fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    run1 = p.add_run()
    run1._r.append(fldChar1)

    instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
    run2 = p.add_run()
    run2._r.append(instrText)

    fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run3 = p.add_run()
    run3._r.append(fldChar2)


def _add_cover_page(doc: Document, meta: dict) -> None:
    """Add a styled cover page."""
    # Add some spacing
    for _ in range(6):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(0)

    # Colored bar
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("━" * 40)
    run.font.color.rgb = BLUE_MED
    run.font.size = Pt(14)

    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(24)
    run = p.add_run(meta.get("title", "Document"))
    run.font.size = Pt(32)
    run.font.color.rgb = BLUE_DARK
    run.bold = True
    run.font.name = "Calibri"

    # Subtitle
    if meta.get("subtitle"):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(meta["subtitle"])
        run.font.size = Pt(16)
        run.font.color.rgb = BLUE_MED
        run.italic = True
        run.font.name = "Calibri"

    # Spacing
    for _ in range(4):
        doc.add_paragraph()

    # Bar
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("━" * 40)
    run.font.color.rgb = BLUE_MED
    run.font.size = Pt(14)

    # Metadata
    info_items = [
        ("Auteur", meta.get("author", "Equipe NordTransit")),
        ("Date", meta.get("date", "2026")),
        ("Version", meta.get("version", "1.0")),
    ]
    for label, value in info_items:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(f"{label} : ")
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_MED
        run.font.name = "Calibri"
        run = p.add_run(value)
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_DARK
        run.bold = True
        run.font.name = "Calibri"

    # Classification
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(36)
    run = p.add_run("CONFIDENTIEL PROJET")
    run.font.size = Pt(10)
    run.font.color.rgb = RED
    run.bold = True

    # Page break
    doc.add_page_break()


def _add_toc(doc: Document) -> None:
    """Add a table of contents field."""
    p = doc.add_heading("Table des matieres", level=1)

    # TOC field
    paragraph = doc.add_paragraph()
    run = paragraph.add_run()
    fldChar = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    run._r.append(fldChar)

    run = paragraph.add_run()
    instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> TOC \\o "1-3" \\h \\z \\u </w:instrText>')
    run._r.append(instrText)

    run = paragraph.add_run()
    fldChar = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="separate"/>')
    run._r.append(fldChar)

    run = paragraph.add_run("[Mettre a jour la table des matieres dans Word : clic droit → Mettre a jour les champs]")
    run.font.color.rgb = GRAY_MED
    run.font.italic = True
    run.font.size = Pt(10)

    run = paragraph.add_run()
    fldChar = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run._r.append(fldChar)

    doc.add_page_break()


def build_docx(md_path: Path) -> Path:
    """Convert a markdown file to a styled DOCX."""
    text = md_path.read_text(encoding="utf-8")
    meta, body = parse_frontmatter(text)
    blocks = parse_markdown_blocks(body)

    doc = Document()

    # Page setup
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2)

    # Default font
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10.5)
    style.font.color.rgb = GRAY_DARK
    style.paragraph_format.space_after = Pt(4)
    style.paragraph_format.line_spacing = 1.15

    # Heading styles
    for level in range(1, 7):
        try:
            hs = doc.styles[f"Heading {level}"]
        except KeyError:
            continue
        hs.font.name = "Calibri"
        hs.font.color.rgb = BLUE_DARK if level <= 2 else BLUE_MED
        hs.font.bold = True
        sizes = {1: 22, 2: 16, 3: 13, 4: 12, 5: 11, 6: 10.5}
        hs.font.size = Pt(sizes.get(level, 11))
        hs.paragraph_format.space_before = Pt(18 if level <= 2 else 12)
        hs.paragraph_format.space_after = Pt(8 if level <= 2 else 4)

    # Cover page
    _add_cover_page(doc, meta)

    # TOC
    _add_toc(doc)

    # Page numbers
    _add_page_number(doc)

    # Render blocks
    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            level = min(block["level"], 6)
            p = doc.add_heading(level=level)
            _add_formatted_runs(p, block["text"])

        elif btype == "paragraph":
            p = doc.add_paragraph()
            _add_formatted_runs(p, block["text"])

        elif btype == "blockquote":
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(1)
            pPr = p._p.get_or_add_pPr()
            pBdr = parse_xml(
                f'<w:pBdr {nsdecls("w")}>'
                '  <w:left w:val="single" w:sz="12" w:space="8" w:color="2E75B6"/>'
                '</w:pBdr>'
            )
            pPr.append(pBdr)
            run = p.add_run(block["text"])
            run.font.italic = True
            run.font.color.rgb = BLUE_MED

        elif btype == "list":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Bullet")
                _add_formatted_runs(p, item)

        elif btype == "numbered_list":
            for item in block["items"]:
                p = doc.add_paragraph(style="List Number")
                _add_formatted_runs(p, item)

        elif btype == "table":
            rows = block["rows"]
            if not rows:
                continue
            ncols = max(len(r) for r in rows)
            table = doc.add_table(rows=len(rows), cols=ncols)
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            table.autofit = True
            _set_table_borders(table)

            for r_idx, row_data in enumerate(rows):
                for c_idx, cell_text in enumerate(row_data):
                    if c_idx >= ncols:
                        break
                    cell = table.cell(r_idx, c_idx)
                    cell.text = ""
                    p = cell.paragraphs[0]
                    p.paragraph_format.space_before = Pt(2)
                    p.paragraph_format.space_after = Pt(2)
                    _add_formatted_runs(p, cell_text)

                    if r_idx == 0:
                        # Header row styling
                        _set_cell_shading(cell, "1B3A5C")
                        for run in p.runs:
                            run.font.color.rgb = WHITE
                            run.font.bold = True
                            run.font.size = Pt(10)
                    else:
                        for run in p.runs:
                            run.font.size = Pt(9.5)
                        # Alternate row coloring
                        if r_idx % 2 == 0:
                            _set_cell_shading(cell, "E8F0F8")

            # Add spacing after table
            doc.add_paragraph()

        elif btype == "code":
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F5F5F5"/>')
            pPr.append(shading)
            pBdr = parse_xml(
                f'<w:pBdr {nsdecls("w")}>'
                '  <w:top w:val="single" w:sz="4" w:space="4" w:color="DDDDDD"/>'
                '  <w:left w:val="single" w:sz="4" w:space="8" w:color="DDDDDD"/>'
                '  <w:bottom w:val="single" w:sz="4" w:space="4" w:color="DDDDDD"/>'
                '  <w:right w:val="single" w:sz="4" w:space="8" w:color="DDDDDD"/>'
                '</w:pBdr>'
            )
            pPr.append(pBdr)
            for code_line in block["text"].split("\n"):
                run = p.add_run(code_line + "\n")
                run.font.name = "Consolas"
                run.font.size = Pt(8.5)
                run.font.color.rgb = GRAY_DARK

        elif btype == "mermaid":
            # Styled placeholder for mermaid diagrams
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pPr = p._p.get_or_add_pPr()
            shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="E8F0F8"/>')
            pPr.append(shading)
            pBdr = parse_xml(
                f'<w:pBdr {nsdecls("w")}>'
                '  <w:top w:val="single" w:sz="8" w:space="8" w:color="2E75B6"/>'
                '  <w:left w:val="single" w:sz="8" w:space="12" w:color="2E75B6"/>'
                '  <w:bottom w:val="single" w:sz="8" w:space="8" w:color="2E75B6"/>'
                '  <w:right w:val="single" w:sz="8" w:space="12" w:color="2E75B6"/>'
                '</w:pBdr>'
            )
            pPr.append(pBdr)

            run = p.add_run("SCHEMA : ")
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = BLUE_DARK
            run = p.add_run(block["description"])
            run.font.size = Pt(10)
            run.font.color.rgb = BLUE_MED
            run.font.italic = True

            # Add a note about the source
            p2 = doc.add_paragraph()
            p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p2.add_run("(Voir le fichier Markdown source pour le diagramme Mermaid interactif)")
            run.font.size = Pt(8)
            run.font.color.rgb = GRAY_MED
            run.font.italic = True

        elif btype == "hr":
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run("━" * 50)
            run.font.color.rgb = BLUE_LIGHT
            run.font.size = Pt(8)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / f"{md_path.stem}.docx"
    doc.save(str(out_path))
    return out_path


# ============================================================
# PPTX BUILDER
# ============================================================

def _pptx_add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = PptxRGB(0x1B, 0x3A, 0x5C)
        run.font.size = PptxPt(36)
        run.font.bold = True
    slide.placeholders[1].text = subtitle
    for run in slide.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = PptxRGB(0x2E, 0x75, 0xB6)
        run.font.size = PptxPt(18)


def _pptx_add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = PptxRGB(0x1B, 0x3A, 0x5C)
        run.font.size = PptxPt(28)
        run.font.bold = True

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        level = 1 if line.startswith("  ") else 0
        p.text = line.strip().lstrip("-•").strip()
        p.level = level
        p.font.size = PptxPt(18 if level == 0 else 16)
        p.font.color.rgb = PptxRGB(0x33, 0x33, 0x33)


def _pptx_add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    txBox = slide.shapes.add_textbox(PptxCm(1.5), PptxCm(0.8), PptxCm(22), PptxCm(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(28)
    p.font.bold = True
    p.font.color.rgb = PptxRGB(0x1B, 0x3A, 0x5C)

    cols = len(headers)
    n_rows = len(rows) + 1
    shape = slide.shapes.add_table(n_rows, cols, PptxCm(1.5), PptxCm(3), PptxCm(22), PptxCm(min(1 * n_rows, 10)))
    tbl = shape.table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = PptxPt(13)
            para.font.bold = True
            para.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptxRGB(0x1B, 0x3A, 0x5C)

    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            if c >= cols:
                break
            cell = tbl.cell(r + 1, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = PptxPt(11)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptxRGB(0xD6, 0xE4, 0xF0)


def build_pptx() -> Path:
    """Build the soutenance PowerPoint."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PptxCm(25.4)
    prs.slide_height = PptxCm(14.29)

    _pptx_add_title_slide(prs,
        "MSPR - NordTransit Logistics",
        "Modernisation infrastructure SI\nEquipe NordTransit | Fevrier 2026")

    _pptx_add_content_slide(prs, "Agenda", [
        "- Contexte et problematique (3-4 min)",
        "- Architecture cible (5-6 min)",
        "- Demo POC live (5 min)",
        "- Migration et budget (3-4 min)",
        "- Conclusion et questions (12 min)",
    ])

    # --- Contexte ---
    _pptx_add_content_slide(prs, "NordTransit Logistics", [
        "- PME logistique — 240 employes, ~100 postes",
        "- 4 entrepots : Lille, Lens, Valenciennes, Arras",
        "- 1 cross-dock saisonnier",
        "- DSI : 4 personnes (1 resp. + 1 admin + 1 tech + 1 alternant)",
    ])

    _pptx_add_content_slide(prs, "Problematique", [
        "- WMS = coeur de metier, 1 serveur, 0 redondance",
        "  Panne = arret total des 4 sites",
        "- Securite insuffisante (MFA partiel, pas de segmentation)",
        "- Pas de PRA (Plan de Reprise d'Activite)",
        "- QoS VoIP non documentee ni configuree",
    ])

    _pptx_add_table_slide(prs, "3 enjeux critiques",
        ["Enjeu", "Impact", "Risque"],
        [
            ["WMS indisponible", "Arret operations 4 sites", "Critique"],
            ["Maintenance reduite", "Nuit uniquement (apres 18h30)", "Eleve"],
            ["DSI de 4 personnes", "Solutions simples obligatoires", "Moyen"],
        ])

    # --- Architecture ---
    _pptx_add_content_slide(prs, "Architecture cible", [
        "- FortiGate (60F/100F) remplace les DrayTek",
        "  VPN site-a-site IKEv2, QoS VoIP, segmentation VLAN",
        "- Cluster 2 serveurs Dell R650xs + SAN PowerVault",
        "  Plus de SPOF serveur",
        "- AD : 3 Domain Controllers (2 on-prem + 1 Azure)",
        "- PRA Azure : Site Recovery pour 20 VMs",
    ])

    _pptx_add_table_slide(prs, "Plan VLAN harmonise",
        ["VLAN", "Nom", "Usage", "QoS"],
        [
            ["10", "MGMT", "Administration", "CS2"],
            ["20", "SERVEURS", "VMs, SAN", "AF31"],
            ["30", "DATA", "Postes, imprimantes", "Best Effort"],
            ["40", "VOIP", "Telephones IP", "EF (prioritaire)"],
        ])

    _pptx_add_content_slide(prs, "Haute disponibilite", [
        "- Cluster 2 noeuds : failover automatique des VMs",
        "- AD multi-DC : DC01 tombe → DC02 prend le relais",
        "- PRA Azure : replication continue, RTO < 1h",
        "- QoS : VoIP prioritaire meme sous charge reseau",
        "- Cloud Witness Azure : quorum sans 3eme serveur",
    ])

    # --- POC ---
    _pptx_add_table_slide(prs, "Lab POC — 7 VMs Proxmox",
        ["VM", "OS", "IP", "Role"],
        [
            ["FW-SIEGE", "pfSense", "172.16.132.1", "Firewall + QoS"],
            ["FW-AZURE", "pfSense", "10.100.0.1", "Cote cloud"],
            ["DC01 + DC02", "Win Server 2022", ".10 / .11", "AD siege"],
            ["DC-AZURE", "Win Server 2022", "10.100.0.10", "AD cloud"],
            ["WMS", "Ubuntu 22.04", "172.16.132.20", "WMS + MySQL"],
            ["IPBX", "FreePBX", "172.16.132.30", "Telephonie VoIP"],
        ])

    _pptx_add_table_slide(prs, "Resultats POC — 4/4 PASS",
        ["Test", "Critere", "Resultat", "Statut"],
        [
            ["QoS VoIP", "Latence < 150ms", "0.1ms, 0% perte", "PASS"],
            ["Failover AD", "Bascule < 15min", "Immediate (< 30s)", "PASS"],
            ["Failover WMS", "Donnees intactes", "5/5 records OK", "PASS"],
            ["Tunnel Azure", "Ping + DNS OK", "< 1ms, replication OK", "PASS"],
        ])

    _pptx_add_content_slide(prs, "Demo live", [
        "- Test QoS : saturation 500Mbps + ping VoIP stable",
        "- Test Failover : arret DC01 → DC02 repond immediatement",
        "- Resultat : architecture validee par les metriques",
    ])

    # --- Migration ---
    _pptx_add_table_slide(prs, "Migration — 6 phases / 6 semaines",
        ["Phase", "Contenu", "Semaine", "Risque"],
        [
            ["M1", "Reseau (FortiGate + VLAN + QoS)", "1", "Moyen"],
            ["M2", "Serveurs (cluster + SAN)", "2", "Eleve"],
            ["M3", "AD (DC02 + replication)", "3", "Faible"],
            ["M4", "WMS (migration cluster)", "4", "Eleve"],
            ["M5", "Azure (PRA + DC cloud)", "5", "Moyen"],
            ["M6", "VoIP (IPBX + QoS)", "6", "Faible"],
        ])

    _pptx_add_table_slide(prs, "Budget — 137k EUR (< 150k)",
        ["Poste", "Montant", "% total"],
        [
            ["Securite (FortiGate x5 + licences)", "17 520 EUR", "13%"],
            ["Serveurs (2x R650xs)", "30 000 EUR", "22%"],
            ["Stockage (SAN ME5012)", "20 000 EUR", "15%"],
            ["Reseau (Cisco C9200 x4)", "12 000 EUR", "9%"],
            ["Azure PRA (3 ans)", "11 232 EUR", "8%"],
            ["Internet + 4G backup", "12 840 EUR", "9%"],
            ["Supervision + Backup", "7 500 EUR", "5%"],
            ["Main d'oeuvre", "8 300 EUR", "6%"],
            ["Marge imprevus (15%)", "17 909 EUR", "13%"],
        ])

    # --- Conclusion ---
    _pptx_add_content_slide(prs, "Avant / Apres", [
        "- SPOF partout → Haute disponibilite (cluster + multi-DC)",
        "- Pas de PRA → PRA Azure (Site Recovery + DC cloud)",
        "- QoS absente → VoIP garantie (prouve par le POC)",
        "- Securite faible → FortiGate + segmentation + VPN IKEv2",
        "- Budget : 137k EUR < 150k EUR (marge 13k)",
    ])

    _pptx_add_content_slide(prs, "Prochaines etapes", [
        "- Formation equipe DSI (4 personnes)",
        "- Phase M1 : deploiement reseau FortiGate",
        "- Monitoring : PRTG + FortiAnalyzer",
        "- Audit securite a 6 mois post-migration",
    ])

    _pptx_add_content_slide(prs, "Merci — Questions ?", [
        "- Documentation complete disponible (6 livrables)",
        "- Lab POC reproductible (7 guides pas-a-pas)",
        "- Budget detaille et justifie poste par poste",
    ])

    # --- Backup Q&A ---
    _pptx_add_table_slide(prs, "Questions anticipees (backup)",
        ["Question", "Reponse cle"],
        [
            ["Pourquoi 2 serveurs, pas 3 ?", "Azure = PRA geo, Cloud Witness = quorum"],
            ["Pourquoi FortiGate ?", "Prix PME, support FR, ecosystem unifie"],
            ["RTO/RPO ?", "AD < 15min/0, WMS < 1h/15min"],
            ["Risques migration ?", "6 phases avec rollback, fenetres nuit"],
            ["Budget realiste ?", "137k < 150k, marge 15% incluse"],
        ])

    out_path = OUTPUT_DIR / "soutenance-nordtransit.pptx"
    prs.save(str(out_path))
    return out_path


# ============================================================
# MAIN
# ============================================================

def main():
    parser = argparse.ArgumentParser(description="Build corporate docs")
    parser.add_argument("--word", action="store_true")
    parser.add_argument("--slides", action="store_true")
    parser.add_argument("--all", action="store_true")
    args = parser.parse_args()

    if not any([args.word, args.slides, args.all]):
        parser.print_help()
        return

    if args.word or args.all:
        print("\n=== Word Livrables (python-docx) ===")
        for md in LIVRABLE_FILES:
            if md.exists():
                out = build_docx(md)
                print(f"  OK: {out.name}")
            else:
                print(f"  SKIP: {md.name}")

    if args.slides or args.all:
        print("\n=== PowerPoint Soutenance ===")
        out = build_pptx()
        print(f"  OK: {out.name}")

    print(f"\nFichiers dans: {OUTPUT_DIR}/")


if __name__ == "__main__":
    main()
